from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import os
import json
from LLM import LLM

class PPTGenerator:
    def __init__(self):
        self.presentation = None
        
    def create_ppt(self, title, content_data, output_path="generated_presentation.pptx"):
        """
        创建PPT演示文稿
        :param title: PPT标题
        :param content_data: 内容数据，格式为列表，每个元素包含slide信息
        :param output_path: 输出文件路径
        :return: 生成结果信息
        """
        try:
            # 创建新的演示文稿
            self.presentation = Presentation()
            
            # 添加标题页
            self._add_title_slide(title)
            
            # 添加内容页
            for slide_data in content_data:
                if slide_data['type'] == 'content':
                    self._add_content_slide(slide_data['title'], slide_data['points'])
                elif slide_data['type'] == 'image':
                    self._add_image_slide(slide_data['title'], slide_data.get('image_path', ''))
                elif slide_data['type'] == 'conclusion':
                    self._add_conclusion_slide(slide_data['title'], slide_data['content'])
            
            # 保存PPT
            full_path = os.path.join(os.getcwd(), output_path)
            self.presentation.save(full_path)
            
            return f"PPT演示文稿已成功生成：{full_path}"
            
        except Exception as e:
            return f"生成PPT时出错：{str(e)}"
    
    def _add_title_slide(self, title):
        """添加标题页"""
        slide_layout = self.presentation.slide_layouts[0]  # 标题页布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        # 设置副标题
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = "由AI助手自动生成"
    
    def _add_content_slide(self, slide_title, bullet_points):
        """添加内容页"""
        slide_layout = self.presentation.slide_layouts[1]  # 标题和内容布局
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_title
        
        # 设置内容
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
    
    def _add_image_slide(self, slide_title, image_path):
        """添加图片页（暂时用文本替代）"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_title
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = f"[图片位置：{image_path}]" if image_path else "[此处应插入相关图片]"
    
    def _add_conclusion_slide(self, slide_title, content):
        """添加总结页"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_title
        
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content

def generate_ppt_from_text(topic: str, requirements: str = "", slide_count: int = 8, style: str = "商务风格"):
    """
    根据主题生成PPT内容结构，使用LLM动态生成内容，确保返回JSON格式
    :param topic: PPT主题
    :param requirements: 特殊要求
    :param slide_count: 幻灯片数量（不含标题页）
    :param style: PPT风格
    :return: PPT内容数据
    """
    max_retries = 3
    for attempt in range(max_retries):
        try:
            # 初始化LLM
            llm = LLM(
                task="你是一个专业的PPT内容生成助手，生成结构化、简洁的幻灯片内容，严格返回JSON格式。",
                model="qwen-plus",
                temperature=0.3,  # 降低temperature以确保结构化输出
            )
            
            # 构造更清晰的提示
            prompt = f"""
            你是一个专业的PPT内容生成助手。请为主题“{topic}”生成PPT内容，包含{slide_count}页（不含标题页），风格为{style}。
            要求：
            - 每页包含一个标题和4-5个要点（content类型幻灯片）或一段总结文本（conclusion类型幻灯片）。
            - 内容必须与“{topic}”高度相关，语言专业，结构清晰，适合{requirements or '通用'}场景。
            - 严格返回以下JSON格式，内容必须包裹在```json```代码块中，确保JSON有效：
            ```json
            [
                {{"type": "content", "title": "概述", "points": ["点1", "点2", "点3", "点4"]}},
                {{"type": "content", "title": "详细分析", "points": ["点1", "点2", "点3", "点4"]}},
                {{"type": "content", "title": "应用场景", "points": ["点1", "点2", "点3", "点4"]}},
                {{"type": "conclusion", "title": "总结与展望", "content": "总结文本"}}
            ]
            ```
            - 不要包含额外文本或注释，仅返回```json```块内的有效JSON。
            示例：
            ```json
            [
                {{"type": "content", "title": "汽车行业概述", "points": ["全球市场规模", "主要厂商", "电动车趋势", "政策影响"]}},
                {{"type": "content", "title": "技术创新", "points": ["自动驾驶技术", "电池技术进步", "车联网发展", "智能制造"]}},
                {{"type": "conclusion", "title": "总结", "content": "汽车行业正迈向电动化与智能化，未来潜力巨大。"}}
            ]
            ```
            """
            
            # 调用LLM生成内容
            response = llm.response(prompt)
            content = response.message.content
            
            # 提取JSON内容
            if content.startswith("```json") and content.endswith("```"):
                json_str = content[7:-3].strip()
            else:
                json_str = content.strip()
                
            content_data = json.loads(json_str)
            
            # 验证内容格式
            if not isinstance(content_data, list):
                raise ValueError("LLM返回的内容不是有效的JSON列表")
            
            # 调整幻灯片数量
            if len(content_data) > slide_count:
                content_data = content_data[:slide_count]
            elif len(content_data) < slide_count:
                for i in range(len(content_data), slide_count):
                    content_data.append({
                        "type": "content",
                        "title": f"补充页{i+1}",
                        "points": [f"{topic}的补充信息{i+1}.1", f"补充信息{i+1}.2", f"补充信息{i+1}.3"]
                    })
            
            return content_data
        
        except json.JSONDecodeError as e:
            print(f"尝试 {attempt+1}/{max_retries}：LLM返回的JSON格式无效 - {str(e)}")
            if attempt == max_retries - 1:
                print("达到最大重试次数，切换到默认模板")
        except Exception as e:
            print(f"尝试 {attempt+1}/{max_retries}：生成PPT内容时出错 - {str(e)}")
            if attempt == max_retries - 1:
                print("达到最大重试次数，切换到默认模板")
    
    # 回退到默认模板
    base_structure = [
        {
            "type": "content",
            "title": "概述",
            "points": [
                f"关于{topic}的基本介绍",
                "主要特点和重要性",
                "当前发展状况",
                "应用领域"
            ]
        },
        {
            "type": "content", 
            "title": "详细分析",
            "points": [
                f"{topic}的核心要素",
                "技术原理或实现方法",
                "优势与挑战",
                "发展趋势"
            ]
        },
        {
            "type": "content",
            "title": "实际应用",
            "points": [
                "典型应用场景",
                "成功案例分析", 
                "实施策略",
                "效果评估"
            ]
        },
        {
            "type": "conclusion",
            "title": "总结与展望",
            "content": f"通过对{topic}的深入分析，我们可以看到其巨大的发展潜力和广阔的应用前景。"
        }
    ]
    
    # 根据slide_count调整
    if slide_count < len(base_structure):
        base_structure = base_structure[:slide_count]
    elif slide_count > len(base_structure):
        for i in range(len(base_structure), slide_count):
            base_structure.append({
                "type": "content",
                "title": f"补充页{i+1}",
                "points": [f"{topic}的补充信息{i+1}.1", f"补充信息{i+1}.2", f"补充信息{i+1}.3"]
            })
    
    # 根据requirements调整
    if requirements and "简短" in requirements:
        base_structure = base_structure[:2]
    elif requirements and "详细" in requirements:
        base_structure.insert(2, {
            "type": "content",
            "title": "深入探讨", 
            "points": [
                f"{topic}的技术细节",
                "相关理论基础",
                "实验数据分析",
                "对比研究"
            ]
        })
    
    return base_structure

def create_presentation(topic, requirements="", filename="", slide_count=8, style="商务风格"):
    """
    为工具调用提供的接口函数
    :param topic: PPT主题
    :param requirements: 特殊要求
    :param filename: 文件名（可选）
    :param slide_count: 幻灯片数量（不含标题页）
    :param style: PPT风格
    :return: 生成结果
    """
    try:
        generator = PPTGenerator()
        content_data = generate_ppt_from_text(topic, requirements, slide_count, style)
        
        # 生成文件名
        if not filename:
            filename = f"{topic}_presentation.pptx"
        if not filename.endswith('.pptx'):
            filename += '.pptx'
            
        result = generator.create_ppt(
            title=f"{topic} - 专题演示",
            content_data=content_data, 
            output_path=filename
        )
        
        return result
        
    except Exception as e:
        return f"创建PPT时出错：{str(e)}"

# 测试功能
if __name__ == "__main__":
    # 测试PPT生成
    result = create_presentation("智能行业概况", "简洁明了", "智能行业PPT", 5, "教学风格")
    print(result)